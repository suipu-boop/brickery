"""渲染底座单测：model -> render 生成文件，断言可重开且为原生元素。

运行：
    python -m pytest tests -v
    （使用已装 python-pptx 的解释器，如 brickery/temp/python/bin/python）
"""

from __future__ import annotations

import os
import sys
import tempfile

# 保证 tests 目录下可 import 本包（tests -> ppt_brick -> brickery 三层）
sys.path.insert(
    0,
    os.path.dirname(
        os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    ),
)

import pytest

from ppt_brick import (
    Box,
    ImageContent,
    ShapeType,
    Slide,
    TextContent,
    render_pptx,
    text_box,
    shape_box,
    HAlign,
)


def _demo_slides():
    slide = Slide(width=13.333, height=7.5, bg="FFFFFF")
    slide.add(
        text_box(
            1.0, 0.8, 10.0, 1.4, "Title", font_size=40.0,
            bold=True, color="1F3864", align=HAlign.CENTER,
        )
    )
    slide.add(shape_box(1.0, 2.4, 10.0, 0.12, fill="2E75B6"))
    slide.add(
        text_box(
            1.0, 2.9, 10.0, 2.5, "Body line 1\nBody line 2",
            font_size=18.0, color="333333",
        )
    )
    return [slide]


def test_render_creates_file(tmp_path):
    out = str(tmp_path / "test.pptx")
    render_pptx(_demo_slides(), out)
    assert os.path.exists(out)
    assert os.path.getsize(out) > 0


def test_reopen_with_python_pptx(tmp_path):
    out = str(tmp_path / "reopen.pptx")
    render_pptx(_demo_slides(), out)

    from pptx import Presentation

    prs = Presentation(out)
    assert len(prs.slides) == 1
    shapes = list(prs.slides[0].shapes)
    # 标题 + 矩形 + 正文 = 3 个原生形状
    assert len(shapes) == 3


def test_unsupported_shape_type_rejected(tmp_path):
    """SVG 复杂特性（如 textPath）不应渲染为可编辑对象，应显式拒绝。"""
    slide = Slide()
    slide.add(
        Box(
            type="shape", x=0.5, y=0.5, w=2, h=2,
            content=__import__(
                "ppt_brick.model", fromlist=["ShapeContent"]
            ).ShapeContent(shape_type="textPath"),
            name="bad",
        )
    )
    with pytest.raises(ValueError):
        render_pptx([slide], str(tmp_path / "x.pptx"))


def test_image_box_renders_native_picture(tmp_path):
    """图片应渲染为原生 picture 对象，而非常景图。"""
    png_path = str(tmp_path / "pixel.png")
    # 1x1 红色 PNG（最小合法 PNG）
    import base64
    png_b64 = (
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mP8z8BQDwAEhQGAhKm"
        "M9QAAAABJRU5ErkJggg=="
    )
    with open(png_path, "wb") as f:
        f.write(base64.b64decode(png_b64))

    slide = Slide()
    slide.add(
        Box(
            type="image", x=0.5, y=0.5, w=2, h=2,
            content=ImageContent(src=png_path),
            name="pic",
        )
    )
    out = str(tmp_path / "img.pptx")
    render_pptx([slide], out)

    from pptx import Presentation

    prs = Presentation(out)
    shapes = list(prs.slides[0].shapes)
    assert len(shapes) == 1
    # 原生 picture（非 shape / placeholder）
    assert shapes[0].shape_type == 13  # PICTURE


if __name__ == "__main__":
    sys.exit(pytest.main([__file__, "-v"]))
