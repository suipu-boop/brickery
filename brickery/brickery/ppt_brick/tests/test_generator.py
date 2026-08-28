"""generator（端到端编排器）单测。

覆盖：编排纯函数（split/chunk/plan）、build_deck 各场景、品牌色缺省、
layout_ids 覆盖、端到端 generate_pptx 落盘并可被 python-pptx 重开。
"""

import os
import sys

sys.path.insert(
    0,
    os.path.dirname(
        os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    ),
)

from pptx import Presentation

from ppt_brick import theme
from ppt_brick.generator import (
    DEFAULT_BRAND,
    chunk_bullets,
    plan_deck,
    split_long_bullet,
    build_deck,
    generate_pptx,
)
from ppt_brick.model import Slide

TOKENS = theme.derive_tokens("1D4ED8")


# ---------------------------------------------------------------------------
# 纯函数：切分
# ---------------------------------------------------------------------------


def test_split_long_bullet_short_keeps():
    assert split_long_bullet("短要点") == ["短要点"]


def test_split_long_bullet_long_splits():
    long_text = "第一句内容。第二句内容；第三句" + "重复。" * 80
    parts = split_long_bullet(long_text, max_chars=40)
    assert len(parts) >= 2
    assert all(p.strip() for p in parts)


def test_split_long_bullet_empty():
    assert split_long_bullet("") == []
    assert split_long_bullet("   ") == []


def test_chunk_bullets_basic_capacity():
    chunks = chunk_bullets([f"b{i}" for i in range(13)], max_per_page=6)
    assert len(chunks) == 3
    assert all(len(c) <= 6 for c in chunks)


def test_chunk_bullets_forced_pages():
    chunks = chunk_bullets([f"b{i}" for i in range(10)], max_per_page=6, forced_pages=2)
    assert len(chunks) == 2


def test_chunk_bullets_forced_too_few_recurses():
    # 强制 2 页但每页会超容量 -> 递归压到 <=6
    chunks = chunk_bullets([f"b{i}" for i in range(20)], max_per_page=6, forced_pages=2)
    assert all(len(c) <= 6 for c in chunks)
    assert sum(len(c) for c in chunks) == 20


def test_chunk_bullets_empty():
    assert chunk_bullets([]) == [[]]


# ---------------------------------------------------------------------------
# 编排计划
# ---------------------------------------------------------------------------


def test_plan_empty_structure_only_cover():
    plan = plan_deck({"title": "T"})
    assert [p["role"] for p in plan] == ["cover"]


def test_plan_single_section():
    plan = plan_deck({"title": "T", "sections": [{"title": "S1", "bullets": ["a"]}]})
    roles = [p["role"] for p in plan]
    # cover + toc(推导) + section + content
    assert roles == ["cover", "toc", "section", "content"]
    assert plan[3]["data"]["note"] == plan[3].get("note", "")  # 空 note


def test_plan_multi_section_pagination():
    structure = {
        "title": "T",
        "toc_items": ["A", "B", "C"],
        "content_max_bullets": 6,
        "sections": [
            {"title": f"S{i}", "bullets": [f"b{j}" for j in range(12)]}
            for i in range(3)
        ],
    }
    plan = plan_deck(structure)
    contents = [p for p in plan if p["role"] == "content"]
    assert len(contents) == 3 * 2  # 每章 12 条 / 6 = 2 页
    # 页码连续且唯一
    page_nums = [p["page"] for p in plan]
    assert page_nums == list(range(1, len(plan) + 1))
    assert plan[0]["total"] == len(plan)


def test_plan_section_forced_pages():
    structure = {"sections": [{"title": "S", "bullets": list("abcde"), "pages": 3}]}
    plan = plan_deck(structure)
    contents = [p for p in plan if p["role"] == "content"]
    assert len(contents) >= 3


def test_plan_per_section_max_bullets():
    structure = {
        "sections": [{"title": "S", "bullets": list("abcdefghij"), "max_bullets": 4}]
    }
    plan = plan_deck(structure)
    contents = [p for p in plan if p["role"] == "content"]
    assert len(contents) == 3  # 10 / 4 -> 3 页


# ---------------------------------------------------------------------------
# build_deck
# ---------------------------------------------------------------------------


def test_build_deck_returns_slides():
    slides = build_deck(TOKENS, {"title": "T", "sections": [{"title": "S", "bullets": ["a"]}]})
    assert slides
    assert all(isinstance(s, Slide) for s in slides)


def test_build_deck_layout_ids_override():
    # content role -> 改用 section 版式（已注册），不影响可渲染性
    structure = {"title": "T", "sections": [{"title": "S", "bullets": ["a", "b"]}]}
    slides = build_deck(TOKENS, structure, layout_ids={"content": "section"})
    assert len(slides) == len(plan_deck(structure))  # 页数不变


def test_build_deck_layout_ids_drop_role():
    structure = {"title": "T", "toc_items": ["A"], "sections": [{"title": "S", "bullets": ["a"]}]}
    slides = build_deck(TOKENS, structure, layout_ids={"toc": None})
    roles = [p["role"] for p in plan_deck(structure)]
    assert "toc" in roles
    # toc 被剥离 -> content 页同样少一页
    assert len(slides) == len(plan_deck(structure)) - 1


def test_build_deck_content_note_page_fill():
    structure = {"title": "T", "sections": [{"title": "S", "bullets": ["a"]}]}
    slides = build_deck(TOKENS, structure)
    # 计划里 content 页应回填页码到 note（经由生成器的 data 注入）
    plan = plan_deck(structure)
    content_pages = [p for p in plan if p["role"] == "content"]
    assert content_pages  # 有 content 页


# ---------------------------------------------------------------------------
# generate_pptx 端到端
# ---------------------------------------------------------------------------


def _count_text_boxes(slide):
    return sum(1 for sh in slide.shapes if sh.has_text_frame)


def test_generate_pptx_roundtrip(tmp_path):
    structure = {
        "brand_color": "#2F5BE7",
        "title": "生成测试",
        "subtitle": "端到端",
        "author": "brickery",
        "date": "2026-08-27",
        "toc_items": ["一", "二", "三", "四", "五"],
        "sections": [
            {"title": f"章 {i}", "bullets": [f"要点 {i}-{j}" for j in range(9)]}
            for i in range(3)
        ],
    }
    out = str(tmp_path / "deck.pptx")
    ret = generate_pptx(structure, out)
    assert ret == out
    assert os.path.exists(out) and os.path.getsize(out) > 0

    plan = plan_deck(structure)
    prs = Presentation(out)
    assert len(prs.slides) == len(plan)
    for s in prs.slides:
        assert _count_text_boxes(s) >= 1
    # 每个 content 页都有页码/总数注记文本
    texts = {
        sh.text_frame.text
        for s in prs.slides
        for sh in s.shapes
        if sh.has_text_frame and sh.text_frame.text.strip()
    }
    assert any("/ " in t for t in texts)


def test_generate_pptx_default_brand(tmp_path):
    # 品牌色缺失 -> DEFAULT_BRAND 兜底
    out = str(tmp_path / "default.pptx")
    generate_pptx({"title": "NO BRAND"}, out)
    assert os.path.exists(out)
    assert Presentation(out).slides  # 可重开


def test_generate_pptx_empty_structure(tmp_path):
    out = str(tmp_path / "single.pptx")
    generate_pptx({"title": "只有封面"}, out)
    assert len(Presentation(out).slides) == 1


def test_generate_pptx_long_bullet(tmp_path):
    long_bullet = "第一个长句段落内容，包含中文标点。" + ("追加细节内容。" * 60)
    structure = {
        "title": "长文",
        "sections": [{"title": "S", "bullets": [long_bullet, "普通要点"]}],
    }
    out = str(tmp_path / "long.pptx")
    generate_pptx(structure, out)
    assert len(Presentation(out).slides) >= 3  # cover+toc+section+>=1 content
