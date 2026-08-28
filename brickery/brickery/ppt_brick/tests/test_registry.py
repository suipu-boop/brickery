"""版式注册表单测：注册表查询 + 每个版式产出可被 render 渲染的合法中间态。

运行（用装了 python-pptx 的解释器）：
    python -m pytest tests/test_registry.py -v
"""
from __future__ import annotations

import os
import sys

sys.path.insert(
    0,
    os.path.dirname(
        os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    ),
)

import pytest

from ppt_brick import theme
from ppt_brick.registry import LAYOUTS, REGISTRY

SAMPLE_DATA = {
    "cover": {
        "title": "OKLCH 主题系统",
        "subtitle": "单品牌色派生完整设计系统",
        "date": "2026-08-27",
        "author": "brickery",
    },
    "toc": {"title": "Agenda", "items": ["主题引擎", "版式注册表", "渲染链路", "示例"]},
    "section": {"number": "03", "title": "版式注册表", "subtitle": "Schema/Component 双层契约"},
    "content": {
        "title": "Token 三级结构",
        "items": ["primitives -> semantic -> variant bundles",
                  "OKLCH 派生 tints/shades + 图表系列色",
                  "WCAG-AA 对比度门禁自动保障",
                  "换 variant 即换整包字体/形语/密度"],
        "note": "说明：内容页要点容量上限 8 条。",
    },
}


def test_registry_has_required_layouts():
    """至少 cover/toc/section/content 三个以上版式。"""
    assert {"cover", "toc", "section", "content"} <= set(LAYOUTS)


def test_get_unknown_raises_keyerror():
    with pytest.raises(KeyError):
        REGISTRY.get("no-such-layout")


def test_duplicate_register_raises():
    with pytest.raises(ValueError):
        REGISTRY.register("cover", lambda t, d=None: [])


def test_schema_has_meta_and_default():
    """每个版式字段都有 meta() 语义与 default() 兜底。"""
    for name in LAYOUTS:
        lay = REGISTRY.get(name)
        meta = lay.meta()
        assert meta, f"{name} 的 schema 不应为空"
        for fn, spec in meta.items():
            assert "meta" in spec, f"{name}/{fn} 缺 meta"
            assert "default" in spec, f"{name}/{fn} 缺 default"


def _pytest_run(name, tokens, out):
    from ppt_brick import render_pptx
    from pptx import Presentation

    data = SAMPLE_DATA[name]
    slides = REGISTRY.get(name).render(tokens, data)
    assert slides, f"{name} 应产出至少 1 页"
    render_pptx(slides, out)
    prs = Presentation(out)
    assert len(prs.slides) == len(slides), f"{name} 页数不一致"
    # 每页应有元素
    for s in prs.slides:
        assert len(list(s.shapes)) > 0


@pytest.mark.parametrize("name", ["cover", "toc", "section", "content"])
def test_layout_renders_light(name, tmp_path):
    """浅色档：每个版式产出可被 render 渲染的合法中间态。"""
    _pytest_run(name, theme.derive_tokens("2F5BE7"),
                str(tmp_path / f"{name}-light.pptx"))


@pytest.mark.parametrize("name", ["cover", "toc", "section", "content"])
def test_layout_renders_dark(name, tmp_path):
    """深色档：同样可渲染（背景深、文字浅）。"""
    _pytest_run(name, theme.derive_tokens("2F5BE7", semantics="dark"),
                str(tmp_path / f"{name}-dark.pptx"))


def test_content_defaults_fill(tmp_path):
    """缺失数据字段用 default 兜底，不抛异常且可渲染。"""
    tokens = theme.derive_tokens("0EA5A5")
    slides = REGISTRY.get("cover").render(tokens, None)
    assert len(slides) == 1
    # 空标题也不该崩溃，且产出的文本盒仍在
    from ppt_brick.model import BoxType

    assert any(b.type == BoxType.TEXT for b in slides[0].boxes)
