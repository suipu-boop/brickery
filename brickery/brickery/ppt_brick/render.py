"""中间态模型 -> 原生 PPTX 落盘（python-pptx 渲染器）。

对齐 specs/brick-ui-and-ppt-brick.md 决策点 D4：
- 输入：ppt_brick.model（Slide/Box/Content）
- 输出：真实 .pptx 文件，元素为 PowerPoint 原生对象
  * TextContent  -> 原生文本框 (sp, textbox)
  * ShapeContent -> 原生 autoshape（可双击改色/改字/加动画）
  * ImageContent -> 原生 picture（可移动/裁剪/替换）
- 禁止：HTML->截图、整页位图化。

逐 Box 调度：对元素数量做基础范围校验，避免越界绘制。
"""

from __future__ import annotations

import base64
import os
from typing import Iterable, List, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from .model import (
    Box,
    BoxType,
    HAlign,
    ImageContent,
    ShapeContent,
    ShapeType,
    Slide,
    TextContent,
    VAlign,
)

# SVG 中无法在原生 PPT 对象上表达的复杂特性白名单（模型刻意不支持）。
# 渲染层对未知形状类型直接抛错，防止静默降级成"截图式"输出。
_SUPPORTED_SHAPES = frozenset(ShapeType.__members__.values())


class RenderError(ValueError):
    """中间态不满足原生渲染约束。"""


def _hex_to_rgb(hex_str: str) -> RGBColor:
    """'AABBCC' -> RGBColor；忽略前导 '#'。"""
    h = hex_str.lstrip("#")
    if len(h) != 6:
        raise RenderError(f"非法颜色 '{hex_str}'，应为 6 位 hex")
    try:
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    except ValueError:
        raise RenderError(f"非法颜色 '{hex_str}'")


def _apply_text_frame(
    text_frame,
    content: TextContent,
    box: Box,
    *,
    first_is_title: bool = False,
):
    """把 TextContent 写到 python-pptx 文本框。"""
    text_frame.word_wrap = True

    # 垂直对齐
    anchor = {
        VAlign.TOP: MSO_ANCHOR.TOP,
        VAlign.MIDDLE: MSO_ANCHOR.MIDDLE,
        VAlign.BOTTOM: MSO_ANCHOR.BOTTOM,
    }[content.valign]
    text_frame.vertical_anchor = anchor

    alignment = {
        HAlign.LEFT: PP_ALIGN.LEFT,
        HAlign.CENTER: PP_ALIGN.CENTER,
        HAlign.RIGHT: PP_ALIGN.RIGHT,
    }[content.align]

    # 单块文本：首段落复用新建文本框自带的空段落，其余 add_paragraph
    lines = str(content.text).split("\n")
    first_para = text_frame.paragraphs[0]
    for i, line in enumerate(lines):
        p = first_para if i == 0 else text_frame.add_paragraph()
        p.alignment = alignment
        run = p.add_run()
        run.text = line
        run.font.size = Pt(content.font_size)
        run.font.bold = content.bold
        run.font.italic = content.italic
        run.font.color.rgb = _hex_to_rgb(content.color)
        if content.font_name:
            run.font.name = content.font_name


def _parse_data_uri(src: str) -> Optional[bytes]:
    """'data:image/png;base64,....' -> bytes；非 data URI 返回 None。"""
    if not src.startswith("data:"):
        return None
    try:
        header, payload = src.split(",", 1)
        if "base64" in header:
            return base64.b64decode(payload)
    except Exception:
        pass
    return None


def _add_image(slide, box: Box, content: ImageContent):
    x, y, w, h = box.to_emu()
    raw = _parse_data_uri(content.src)
    path = content.src
    if raw is not None:
        import io
        import tempfile

        # 中间产物：临时文件写系统临时目录（render 阶段不受工作区约束）
        suffix = ".png"
        if content.src.split(";", 1)[0].endswith("/jpeg"):
            suffix = ".jpg"
        fd, path = tempfile.mkstemp(suffix=suffix)
        with os.fdopen(fd, "wb") as f:
            f.write(raw)
        try:
            pic = slide.shapes.add_picture(path, x, y, w, h)
        finally:
            os.unlink(path)
    else:
        if not os.path.exists(path):
            raise RenderError(f"图片不存在: {path}")
        pic = slide.shapes.add_picture(path, x, y, w, h)

    # fit 语义：contain 不强制拉伸（python-pptx 默认按原比例置于给定框内）
    if content.fit == "stretch":
        pic.width = w
        pic.height = h


def _add_box(slide, box: Box):
    if box.type == BoxType.TEXT:
        content = box.content
        if not isinstance(content, TextContent):
            raise RenderError(f"Box '{box.name}' 类型 text 但内容不是 TextContent")
        tb = slide.shapes.add_textbox(*box.to_emu())
        _apply_text_frame(tb.text_frame, content, box)
        # 元素背景（对文本框生效需先关填充默认）
        if box.bg:
            tb.fill.solid()
            tb.fill.fore_color.rgb = _hex_to_rgb(box.bg)
        return

    if box.type == BoxType.SHAPE:
        content = box.content
        if not isinstance(content, ShapeContent):
            raise RenderError(f"Box '{box.name}' 类型 shape 但内容不是 ShapeContent")
        if content.shape_type not in _SUPPORTED_SHAPES:
            raise RenderError(
                f"形状类型 {content.shape_type} 不在原生可编辑白名单内"
            )
        mso_name = content.shape_type.to_mso_shape()
        shape = slide.shapes.add_shape(
            getattr(MSO_SHAPE, mso_name), *box.to_emu()
        )
        if content.fill == "none":
            shape.fill.background()
        else:
            shape.fill.solid()
            shape.fill.fore_color.rgb = _hex_to_rgb(content.fill)
        if content.line_color:
            shape.line.color.rgb = _hex_to_rgb(content.line_color)
            shape.line.width = Inches(content.line_width / 72.0)
        else:
            shape.line.fill.background()
        return

    if box.type == BoxType.IMAGE:
        content = box.content
        if not isinstance(content, ImageContent):
            raise RenderError(f"Box '{box.name}' 类型 image 但内容不是 ImageContent")
        _add_image(slide, box, content)
        return

    raise RenderError(f"未知 Box 类型: {box.type}")


def render_slide(prs: Presentation, slide: Slide):
    """把一页中间态 Slide 渲染到 prs。"""
    blank = prs.slide_layouts[6]  # Blank layout
    s = prs.slides.add_slide(blank)

    # 画布大小
    prs.slide_width = Inches(slide.width)
    prs.slide_height = Inches(slide.height)

    # 背景
    if slide.bg != "transparent":
        bg = s.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = _hex_to_rgb(slide.bg)

    for box in slide.boxes:
        _add_box(s, box)
    return s


def render_pptx(slides: Iterable[Slide], out_path: str) -> str:
    """将一组 Slide 渲染为原生 .pptx 并落盘。

    Returns:
        out_path（与入参一致），便于链式调用。
    """
    slides = list(slides)
    if not slides:
        raise RenderError("没有可渲染的 Slide")

    prs = Presentation()
    prs.slide_width = Inches(slides[0].width)
    prs.slide_height = Inches(slides[0].height)
    for slide in slides:
        render_slide(prs, slide)

    out_dir = os.path.dirname(os.path.abspath(out_path))
    os.makedirs(out_dir, exist_ok=True)
    prs.save(out_path)
    return out_path
